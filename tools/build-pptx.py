#!/usr/bin/env python3
"""build-pptx.py — generate the UI Button Guide as a PowerPoint deck.

Mirrors docs/slides/index.html (light style). Output can be imported
directly into Google Slides (File > Import) or uploaded to Drive.

Usage: python3 tools/build-pptx.py [out.pptx]
"""
import io
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(ROOT, 'docs', 'slides', 'img')
OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(ROOT, 'docs', 'slides', 'VectorScope_UI_Button_Guide.pptx')

INK = RGBColor(0x1F, 0x24, 0x30)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
BG = RGBColor(0xF7, 0xF8, 0xFA)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

# ── deck content ─────────────────────────────────────────────────────
# Each: (title, chip, description, [(tag, image_basename), ...])
SLIDES = [
    ("Screen Layout", "orientation",
     "Five panels: Bird's Eye (top-left), Main / UW / Tele camera views (bottom row), and the Combined "
     "view (bottom-right) — the simulated phone output that switches cameras as you zoom. The Controls "
     "panel (top-right) hosts every button in this deck. Bottom-right shows the live 3\u00d73 sampling matrix.",
     [("Full app", "00_overview"), ("Controls panel", "00_controls_panel")]),

    ("Set Camera", "dialog",
     "Opens the camera parameter editor: intrinsics (fx fy cx cy), extrinsics (position, rotation) for all "
     "three cameras, plus the scene camera pose. Edits apply live. A camera JSON file can also be imported.",
     [("Before", "01_setcam_before"), ("After click", "01_setcam_after")]),

    ("Save Scene \u00b7 Load Scene \u00b7 Reset All", "file ops",
     "Save Scene downloads the full state (objects, cameras, sliders, segment config, curves) as JSON. "
     "Load Scene restores it. Reset All returns every control to defaults — below, zoom was at 4.0x with "
     "Focus D 1.2 m; after Reset All the app is back at 1.0x / 3.0 m.",
     [("Action buttons", "20_file_buttons"), ("Before reset (4.0x, D=1.2)", "18_reset_before"),
      ("After Reset All (1.0x, D=3.0)", "18_reset_after")]),

    ("Help", "dialog",
     "Opens the built-in manual: pipeline explanation, per-segment table, blending modes, radial blend "
     "behavior, trajectory recording, and file formats. Rendered from the project docs.",
     [("Before", "02_help_before"), ("After click", "02_help_after")]),

    ("Focus D slider + AF button", "focus plane",
     "Focus D sets the focus-plane depth used by the plane-induced homographies — objects at this depth "
     "align perfectly between cameras; off-plane objects show parallax. AF arms autofocus: drag a box on "
     "any camera panel and the median depth inside becomes the new Focus D.",
     [("D = 3.0 m (warp off, 0.7x)", "15_focusd_3"), ("D = 1.0 m — alignment shifts", "15_focusd_1"),
      ("AF: box dragged on Main panel", "14_af_after")]),

    ("Prewarp 1 / Prewarp 2 sliders", "warp-off path",
     "When Warp is OFF, camera handovers use a fixed prewarp scale instead of homography interpolation. "
     "Prewarp 1 applies to the UW segment (0.5\u20131.0x), Prewarp 2 to the Main\u2192Tele segment "
     "(2.0\u20135.0x). Below at 0.7x, raising Prewarp 1 from 2.0 \u2192 4.0 changes the Combined framing.",
     [("Prewarp1 = 2.0", "16_prewarp_before"), ("Prewarp1 = 4.0", "16_prewarp_after")]),

    ("Zoom slider + Go presets", "0.5x \u2013 10x",
     "The Zoom slider is log-scaled over 0.5x\u201310x. The Go presets (0.5x / 1x / 2x / 5x / 10x) animate "
     "smoothly to the target using the zoom transition curve. The Combined view's lead camera follows the "
     "segment config: UW below 1x \u2192 Main \u2192 Tele at 5x+.",
     [("1.0x — lead: Main", "11_zoom_1x"), ("Go 5.0x — lead: Tele", "11_zoom_5x"),
      ("Go 0.5x — lead: UW", "11_zoom_05x")]),

    ("Seg — Segment Config", "dialog",
     "Configures the zoom pipeline via breakpoints (defaults: 1.0, 2.0, 5.0) dividing 0.5x\u201310x into "
     "segments. Each segment assigns a Lead camera (shown in Combined), a Follower (used only during "
     "transition blending), and a per-segment Warp flag. Breakpoints auto-sort; below = z < bp, at/above = z \u2265 bp.",
     [("Before", "03_seg_before"), ("After click — segment table", "03_seg_after")]),

    ("\u25d1 Zoom Curve", "bezier editor",
     "Photoshop-style cubic bezier editor for the zoom transition speed used by the Go presets — plus the "
     "transition duration (ms). Drag the four control points to make transitions ease-in, ease-out, or snappy.",
     [("Before", "04_zoomcurve_before"), ("After click — curve editor", "04_zoomcurve_after")]),

    ("Warp + \u25d1 Warp Curve", "homography interpolation",
     "Warp ON: within warp-enabled segments the sampling matrix interpolates lerp(I \u2192 H\u207b\u00b9) on "
     "a log-zoom scale — the view morphs continuously between camera perspectives. Warp OFF: a fixed "
     "prewarp \u00d7 crop is used instead. The \u25d1 button opens a bezier editor shaping the interpolation t.",
     [("Warp ON (0.7x)", "06_warp_on"), ("Warp OFF (0.7x)", "06_warp_off"),
      ("\u25d1 Warp curve editor", "05_warpcurve_after")]),

    ("\u25b6 Play", "auto sweep",
     "Runs an automatic zoom bounce loop: zoom sweeps up to 10x, reverses down to 0.5x, and repeats — a "
     "hands-free demo of every camera handover, warp interpolation, and blend. Click again (\u25a0 Stop) to end.",
     [("Before (1.0x, idle)", "10_play_before"), ("Playing — mid-sweep", "10_play_after")]),

    ("30 / 60 FPS", "render rate",
     "Toggles the fixed render rate between 30 and 60 FPS. Late frames are dropped, never queued — "
     "simulating a real camera pipeline's frame pacing. The button label shows the active rate.",
     [("30 FPS", "09_fps_before"), ("60 FPS", "09_fps_after")]),

    ("Combined", "focus mode",
     "Focus mode: hides all other panels and maximizes the Combined view — the simulated phone output. "
     "Ideal for judging handover quality and blend artifacts at full size. Click again to restore the layout.",
     [("Before — 5 panels", "08_combined_before"), ("After — Combined only", "08_combined_after")]),

    ("Grid", "sampling overlay",
     "Overlays the lead and follower sampling grids on the Combined view — each grid drawn in its camera's "
     "color (UW green, Main blue, Tele yellow). Shows exactly how the warp maps output pixels back into "
     "each camera's frame. Labels appear top-right.",
     [("Before (3.0x)", "07_grid_before"), ("After — grids on Combined", "07_grid_after")]),

    ("Blend slider \u00b7 Single/Dual \u00b7 Flat/Radial", "camera transition",
     "On a lead-camera switch the old and new cameras cross-fade over Blend frames. Single freezes the "
     "outgoing frame; Dual re-renders the follower live every frame. Flat = uniform alpha; Radial = circular "
     "sweep — for Tele\u2192Main the new camera appears from the edges and the Tele region shrinks to the center.",
     [("At 5.0x (Tele), Dual+Radial armed", "13_blend_at5x"), ("Mid-blend \u2192 1.0x — radial sweep", "13_blend_mid"),
      ("Blend complete (Main)", "13_blend_done")]),

    ("Trajectory: \u25cf Rec \u00b7 \u25b6 \u00b7 \u23ee \u23ed \u00b7 \u25a0 \u00b7 Load / Save", "record & replay",
     "\u25cf Rec records every per-frame state (zoom, focus, blend\u2026) into a trajectory. Stop \u2192 it "
     "appears in the dropdown and the transport unlocks: \u25b6 play/pause, \u23ee/\u23ed frame-step, progress "
     "slider to seek, \u25a0 back to free mode. Load/Save exchange trajectory JSON files.",
     [("Recording (zoom moves captured)", "19_traj_recording"), ("Recorded — transport unlocked", "19_traj_recorded"),
      ("Playing back", "19_traj_playing")]),

    ("Ghost Y slider", "bird's eye",
     "Sets the ghost height (world Y) for the Bird's Eye panel: everything above this height renders "
     "translucent so it doesn't block the top-down view. At 2.0 m the ceiling is ghosted and you see inside "
     "the room; at 6.0 m nothing is ghosted, so the solid roof hides the interior.",
     [("Ghost Y = 2.0 — room visible", "17_clip_before"), ("Ghost Y = 6.0 — roof covers view", "17_clip_after")]),
]


def webp_to_png_stream(basename):
    """python-pptx can't embed webp — convert in memory (JPEG keeps size sane)."""
    im = Image.open(os.path.join(IMG, basename + '.webp')).convert('RGB')
    buf = io.BytesIO()
    im.save(buf, 'JPEG', quality=88)
    buf.seek(0)
    return buf, im.size


def add_bg(slide, prs):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ── Title slide ──
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_text(s, Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.9),
             "VectorScope — UI Button Guide", 36, INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.6),
             "Every control explained, with before / after screenshots.  "
             "UW \u00b7 Main \u00b7 Tele — continuous 0.5x\u201310x zoom pipeline.",
             14, MUTED, align=PP_ALIGN.CENTER)
    buf, (iw, ih) = webp_to_png_stream("00_overview")
    pic_w = Inches(9.6)
    pic_h = Emu(int(pic_w * ih / iw))
    s.shapes.add_picture(buf, (SLIDE_W - pic_w) // 2, Inches(2.05), width=pic_w)

    # ── Content slides ──
    for title, chip, desc, shots in SLIDES:
        s = prs.slides.add_slide(blank)
        add_bg(s, prs)
        add_text(s, Inches(0.6), Inches(0.35), Inches(9.5), Inches(0.6), title, 26, INK, bold=True)
        add_text(s, Inches(10.2), Inches(0.45), Inches(2.5), Inches(0.5),
                 chip.upper(), 11, ACCENT, bold=True, align=PP_ALIGN.RIGHT)
        # accent rule
        rule = s.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(12.1), Pt(2.4))
        rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background()
        add_text(s, Inches(0.6), Inches(1.15), Inches(12.1), Inches(1.0), desc, 12.5, MUTED)

        n = len(shots)
        gap = Inches(0.3)
        total_w = SLIDE_W - Inches(1.2) - gap * (n - 1)
        cell_w = Emu(int(total_w / n))
        y_img = Inches(2.45)
        max_h = SLIDE_H - y_img - Inches(0.35)
        for i, (tag, base) in enumerate(shots):
            x = Inches(0.6) + Emu(int((cell_w + gap) * i))
            add_text(s, x, y_img - Inches(0.32), cell_w, Inches(0.3), tag, 10.5, MUTED, bold=True)
            buf, (iw, ih) = webp_to_png_stream(base)
            w, h = cell_w, Emu(int(cell_w * ih / iw))
            if h > max_h:
                h = Emu(int(max_h))
                w = Emu(int(h * iw / ih))
            s.shapes.add_picture(buf, x, y_img, width=w, height=h)

    # ── End slide ──
    s = prs.slides.add_slide(blank)
    add_bg(s, prs)
    add_text(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(0.9),
             "That's every button.", 34, INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.2),
             "Docs: HOMOGRAPHY_PIPELINE.md \u00b7 RADIAL_BLENDING.md \u00b7 MODULE_MAP.md \u00b7 "
             "FILE_FORMATS.md \u00b7 CAMERAS.md\n"
             "Screenshots auto-generated by tools/capture-slides.js \u00b7 deck by tools/build-pptx.py",
             13, MUTED, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print("saved:", OUT, f"({os.path.getsize(OUT) / 1e6:.1f} MB, {len(prs.slides.slides if hasattr(prs.slides, 'slides') else prs.slides._sldIdLst)} slides)")


if __name__ == '__main__':
    main()
